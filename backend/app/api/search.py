"""
Search API endpoints
"""
from fastapi import APIRouter, Depends, Query, Body, HTTPException
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from sqlalchemy import or_, and_
from typing import List, Optional
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os
import re
import tempfile
from datetime import datetime

from app.core.database import get_db
from app.models.slide import Slide
from app.models.tag import Tag

router = APIRouter()


class ExportRequest(BaseModel):
    slide_ids: List[int]
    filename: Optional[str] = None  # Optional custom filename


@router.get("/")
async def search_slides(
    q: Optional[str] = Query(None, description="Search query for text content"),
    tags: Optional[List[str]] = Query(None, description="Tag names to filter by"),
    categories: Optional[List[str]] = Query(None, description="Tag categories to filter by"),
    db: Session = Depends(get_db)
):
    """
    Search slides by text content and/or tags
    
    Examples:
    - /api/search?q=analytics
    - /api/search?tags=enterprise&tags=CRM
    - /api/search?categories=client_type&categories=software_type
    - /api/search?q=dashboard&tags=reporting
    """
    query = db.query(Slide)
    
    # Text search
    if q:
        search_filter = or_(
            Slide.text_content.ilike(f"%{q}%"),
            Slide.title.ilike(f"%{q}%")
        )
        query = query.filter(search_filter)
    
    # Filter by tags
    if tags:
        for tag_name in tags:
            query = query.join(Slide.tags).filter(Tag.name == tag_name)
    
    # Filter by categories
    if categories:
        query = query.join(Slide.tags).filter(Tag.category.in_(categories))
    
    # Execute query
    slides = query.distinct().all()
    
    return {
        "count": len(slides),
        "slides": [
            {
                "id": slide.id,
                "original_filename": slide.original_filename,
                "slide_number": slide.slide_number,
                "thumbnail_path": slide.thumbnail_path,
                "image_path": slide.image_path,
                "title": slide.title,
                "text_content": slide.text_content[:200] if slide.text_content else None,  # Preview
                "tags": [
                    {"id": tag.id, "name": tag.name, "category": tag.category}
                    for tag in slide.tags
                ]
            }
            for slide in slides
        ]
    }


@router.post("/export")
async def export_slides(
    request: ExportRequest,
    db: Session = Depends(get_db)
):
    """
    Export selected slides as a new PowerPoint presentation
    Creates a PPTX file with the selected slides as images
    """
    # Get slides from database
    slides = db.query(Slide).filter(Slide.id.in_(request.slide_ids)).order_by(Slide.id).all()
    
    if not slides:
        raise HTTPException(status_code=404, detail="No slides found with the provided IDs")
    
    try:
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Set slide size to standard 16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add each slide as an image
        for slide in slides:
            # Use the full-size image path
            image_path = slide.image_path
            
            # Check if image exists
            if not os.path.exists(image_path):
                print(f"Warning: Image not found for slide {slide.id}: {image_path}")
                continue
            
            # Add a blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide_obj = prs.slides.add_slide(blank_slide_layout)
            
            # Get image dimensions
            img = Image.open(image_path)
            img_width, img_height = img.size
            img.close()
            
            # Calculate dimensions to fit slide while maintaining aspect ratio
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Calculate scaling
            width_ratio = slide_width / img_width
            height_ratio = slide_height / img_height
            scale_ratio = min(width_ratio, height_ratio)
            
            # Calculate final dimensions
            final_width = int(img_width * scale_ratio)
            final_height = int(img_height * scale_ratio)
            
            # Center the image
            left = (slide_width - final_width) // 2
            top = (slide_height - final_height) // 2
            
            # Add image to slide
            slide_obj.shapes.add_picture(
                image_path,
                left,
                top,
                width=final_width,
                height=final_height
            )
        
        # Determine filename
        if request.filename:
            # Sanitize the filename to remove invalid characters
            safe_filename = re.sub(r'[<>:"/\\|?*]', '_', request.filename)
            # Ensure it ends with .pptx
            if not safe_filename.lower().endswith('.pptx'):
                safe_filename += '.pptx'
            download_filename = safe_filename
        else:
            # Use default timestamped filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            download_filename = f"slidex_export_{timestamp}.pptx"
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".pptx"
        )
        temp_file.close()
        
        prs.save(temp_file.name)
        
        # Return the file
        return FileResponse(
            path=temp_file.name,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=download_filename,
            background=None  # File will be deleted after sending
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error creating PowerPoint: {str(e)}")

# Made with Bob
