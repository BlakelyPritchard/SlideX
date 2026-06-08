"""
Deck Assembler Service for SlideX
Assembles PowerPoint presentations from selected slides with AI-generated presenter notes
"""

import os
import logging
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from app.models.slide import Slide
from app.services.ai_tagger import get_ai_provider
from datetime import datetime

logger = logging.getLogger(__name__)


class DeckAssembler:
    """Service for assembling PowerPoint decks from slides"""
    
    def __init__(self):
        """Initialize the deck assembler"""
        self.ai_provider = get_ai_provider()
        self.slides_dir = os.getenv("SLIDES_DIR", "./slides")
    
    def create_presentation(
        self,
        slides: List[Slide],
        title: str = "Generated Presentation",
        add_agenda: bool = True,
        add_notes: bool = True,
        client_name: Optional[str] = None
    ) -> str:
        """
        Create a PowerPoint presentation from selected slides
        
        Args:
            slides: List of Slide objects to include
            title: Title for the presentation
            add_agenda: Whether to add an agenda slide
            add_notes: Whether to generate presenter notes
            client_name: Optional client name for customization
            
        Returns:
            Path to the created presentation file
        """
        try:
            logger.info(f"Creating presentation with {len(slides)} slides")
            
            # Create new presentation
            prs = Presentation()
            
            # Set slide size to standard 16:9
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Add title slide
            self._add_title_slide(prs, title, client_name)
            
            # Add agenda slide if requested
            if add_agenda and len(slides) > 0:
                self._add_agenda_slide(prs, slides)
            
            # Add content slides
            for i, slide in enumerate(slides):
                logger.info(f"Adding slide {i+1}/{len(slides)}: {slide.title}")
                self._add_content_slide(prs, slide, add_notes)
            
            # Add closing slide
            self._add_closing_slide(prs, client_name)
            
            # Save presentation
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"generated_deck_{timestamp}.pptx"
            output_path = os.path.join(self.slides_dir, filename)
            
            prs.save(output_path)
            logger.info(f"Presentation saved to {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            raise
    
    def _add_title_slide(
        self,
        prs: Presentation,
        title: str,
        client_name: Optional[str]
    ):
        """Add title slide to presentation"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set subtitle with client name if provided
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            if client_name:
                subtitle.text = f"Prepared for {client_name}\n{datetime.now().strftime('%B %Y')}"
            else:
                subtitle.text = datetime.now().strftime('%B %Y')
    
    def _add_agenda_slide(self, prs: Presentation, slides: List[Slide]):
        """Add agenda slide listing main topics"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Agenda"
        
        # Add agenda items
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # Group slides by category/theme
            agenda_items = self._generate_agenda_items(slides)
            
            for item in agenda_items:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(18)
    
    def _generate_agenda_items(self, slides: List[Slide]) -> List[str]:
        """Generate agenda items from slides"""
        # Simple approach: use slide titles or generate themes
        items = []
        seen_themes = set()
        
        for slide in slides[:10]:  # Limit to first 10 for agenda
            if slide.title and slide.title not in seen_themes:
                # Clean up title
                title = slide.title.strip()
                if len(title) > 50:
                    title = title[:47] + "..."
                items.append(title)
                seen_themes.add(slide.title)
        
        # If we have too few items, add generic ones
        if len(items) < 3:
            items = [
                "Introduction & Overview",
                "Key Capabilities",
                "Value Proposition",
                "Next Steps"
            ]
        
        return items
    
    def _add_content_slide(
        self,
        prs: Presentation,
        slide: Slide,
        add_notes: bool
    ):
        """Add a content slide from the repository"""
        try:
            # Create blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            new_slide = prs.slides.add_slide(blank_layout)
            
            # Add slide image - use the full path from the database
            image_path = slide.image_path
            if os.path.exists(image_path):
                # Add image to fill the slide
                left = Inches(0)
                top = Inches(0)
                width = prs.slide_width
                height = prs.slide_height
                
                new_slide.shapes.add_picture(
                    image_path,
                    left, top,
                    width=width,
                    height=height
                )
            else:
                logger.warning(f"Image not found: {image_path}")
                # Add title as fallback
                title_box = new_slide.shapes.add_textbox(
                    Inches(1), Inches(1),
                    Inches(8), Inches(1)
                )
                title_frame = title_box.text_frame
                title_frame.text = slide.title or "Slide"
            
            # Add presenter notes if requested
            if add_notes:
                notes = self._generate_presenter_notes(slide)
                if notes:
                    notes_slide = new_slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = notes
            
        except Exception as e:
            logger.error(f"Error adding content slide: {e}")
            # Continue with other slides even if one fails
    
    def _generate_presenter_notes(self, slide: Slide) -> str:
        """Generate AI-powered presenter notes for a slide"""
        try:
            # Create prompt for AI
            prompt = f"""Generate concise presenter notes for this slide.
Include:
- Key talking points (2-3 bullet points)
- Important statistics or facts to mention
- Transition to next topic

Slide Title: {slide.title or 'Untitled'}
Slide Content: {slide.text_content or 'No text content'}

Keep notes under 150 words. Be specific and actionable."""

            # Generate notes using AI
            notes = self.ai_provider.generate_text(prompt, max_tokens=200)
            
            if notes:
                return notes.strip()
            else:
                return self._generate_fallback_notes(slide)
                
        except Exception as e:
            logger.error(f"Error generating presenter notes: {e}")
            return self._generate_fallback_notes(slide)
    
    def _generate_fallback_notes(self, slide: Slide) -> str:
        """Generate basic presenter notes without AI"""
        notes_parts = []
        
        if slide.title:
            notes_parts.append(f"Topic: {slide.title}")
        
        if slide.text_content:
            # Extract first few sentences
            content = slide.text_content[:200]
            notes_parts.append(f"\nKey Points:\n{content}")
        
        # Add tag-based talking points
        if slide.tags:
            tag_names = [tag.name.replace('_', ' ').title() for tag in slide.tags[:3]]
            notes_parts.append(f"\nTopics: {', '.join(tag_names)}")
        
        return "\n".join(notes_parts) if notes_parts else "Discuss this slide with the audience."
    
    def _add_closing_slide(self, prs: Presentation, client_name: Optional[str]):
        """Add closing slide with next steps"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Next Steps"
        
        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            next_steps = [
                "Schedule a detailed technical discussion",
                "Arrange a proof of concept",
                "Review pricing and licensing options",
                "Connect with our implementation team"
            ]
            
            for step in next_steps:
                p = tf.add_paragraph()
                p.text = step
                p.level = 0
                p.font.size = Pt(18)
    
    def add_transition_slide(
        self,
        prs: Presentation,
        section_title: str,
        description: Optional[str] = None
    ):
        """Add a transition slide between sections"""
        slide_layout = prs.slide_layouts[2]  # Section header layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = section_title
        
        # Add description if provided
        if description and len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            content.text = description
    
    def customize_branding(
        self,
        prs: Presentation,
        logo_path: Optional[str] = None,
        color_scheme: Optional[Dict[str, str]] = None
    ):
        """
        Customize presentation with client branding
        
        Args:
            prs: Presentation object
            logo_path: Path to client logo image
            color_scheme: Dictionary with color values
        """
        # This is a placeholder for future branding customization
        # Would add logo to master slides, update color scheme, etc.
        pass


# Global instance
_deck_assembler: Optional[DeckAssembler] = None


def get_deck_assembler() -> DeckAssembler:
    """
    Get or create the global deck assembler instance
    
    Returns:
        DeckAssembler instance
    """
    global _deck_assembler
    if _deck_assembler is None:
        _deck_assembler = DeckAssembler()
    return _deck_assembler

# Made with Bob
