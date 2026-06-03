"""
Slide processing service
Handles PowerPoint file extraction, image conversion, and text extraction
"""
from pptx import Presentation
from PIL import Image
import io
import os
from typing import List, Dict
from app.core.config import settings


class SlideProcessor:
    """Process PowerPoint presentations into individual slides"""
    
    def __init__(self):
        self.slides_dir = settings.SLIDES_DIR
        self.thumbnails_dir = settings.THUMBNAILS_DIR
    
    async def process_presentation(self, file_path: str, filename: str) -> List[Dict]:
        """
        Process a PowerPoint presentation
        Returns list of slide data dictionaries
        """
        prs = Presentation(file_path)
        slides_data = []
        
        base_filename = os.path.splitext(filename)[0]
        
        for idx, slide in enumerate(prs.slides, start=1):
            # Extract text content
            text_content = self._extract_text_from_slide(slide)
            title = self._extract_title_from_slide(slide)
            
            # Convert slide to image
            image_filename = f"{base_filename}_slide_{idx}.png"
            thumbnail_filename = f"{base_filename}_slide_{idx}_thumb.png"
            
            image_path = os.path.join(self.slides_dir, image_filename)
            thumbnail_path = os.path.join(self.thumbnails_dir, thumbnail_filename)
            
            # Note: python-pptx doesn't directly convert to images
            # For MVP, we'll create placeholder images
            # In production, use libraries like pdf2image or aspose.slides
            self._create_placeholder_image(image_path, (1920, 1080))
            self._create_placeholder_image(thumbnail_path, (320, 240))
            
            slides_data.append({
                "slide_number": idx,
                "image_path": image_path,
                "thumbnail_path": thumbnail_path,
                "text_content": text_content,
                "title": title
            })
        
        return slides_data
    
    def _extract_text_from_slide(self, slide) -> str:
        """Extract all text from a slide"""
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
        return "\n".join(text_runs)
    
    def _extract_title_from_slide(self, slide) -> str:
        """Extract title from slide if it exists"""
        if slide.shapes.title:
            return slide.shapes.title.text
        return ""
    
    def _create_placeholder_image(self, path: str, size: tuple):
        """
        Create a placeholder image
        TODO: Replace with actual slide rendering in production
        """
        img = Image.new('RGB', size, color=(240, 240, 240))
        img.save(path)

# Made with Bob
